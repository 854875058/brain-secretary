from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "docs" / "presentations"
OUTPUT_PPTX = OUTPUT_DIR / "spec-coding-three-agent-share.pptx"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(20, 43, 74)
TEAL = RGBColor(17, 138, 178)
CORAL = RGBColor(242, 95, 92)
GOLD = RGBColor(247, 184, 1)
SLATE = RGBColor(84, 98, 111)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(32, 41, 50)
MUTED = RGBColor(102, 112, 122)

TITLE_FONT = "Microsoft YaHei"
BODY_FONT = "Microsoft YaHei"


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, color: RGBColor = NAVY) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, text: str = "Spec Coding × 三节点 Agent × Brain Secretary") -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(6.6), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.8), Inches(1.1))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = TITLE_FONT
    run.font.size = Pt(27)
    run.font.bold = True
    run.font.color.rgb = WHITE if dark else NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.52), Inches(11.2), Inches(0.5))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = BODY_FONT
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(220, 228, 236) if dark else SLATE


def add_bullets(slide, bullets: list[str], left: float = 0.9, top: float = 2.0, width: float = 11.2, height: float = 4.4) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(12)
        p.font.name = BODY_FONT
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT


def add_two_column(slide, left_title: str, left_items: list[str], right_title: str, right_items: list[str]) -> None:
    left_rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.9), Inches(5.8), Inches(4.6))
    right_rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.9), Inches(5.8), Inches(4.6))
    for rect, fill in ((left_rect, LIGHT), (right_rect, LIGHT)):
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = RGBColor(220, 226, 232)
        rect.line.width = Pt(1.2)

    for x, title, items in ((0.95, left_title, left_items), (7.1, right_title, right_items)):
        title_box = slide.shapes.add_textbox(Inches(x), Inches(2.15), Inches(5.0), Inches(0.4))
        p = title_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = TITLE_FONT
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = NAVY
        add_bullets(slide, items, left=x, top=2.7, width=5.0, height=3.4)


def add_flow_box(slide, x: float, y: float, w: float, h: float, title: str, body: str, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = TITLE_FONT
    r1.font.bold = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = BODY_FONT
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = SLATE) -> None:
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(2.2)


def build_deck() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title(slide, "Spec Coding 项目实战", "从单 AI 辅助到三节点 Agent 交付闭环\n以 Brain Secretary 为例", dark=True)
    quote = slide.shapes.add_textbox(Inches(0.8), Inches(2.45), Inches(11.8), Inches(1.6))
    p = quote.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "核心观点：Spec 不是长 Prompt，而是把需求、边界、验证和输出格式写清楚，再交给多节点 Agent 去执行。"
    r.font.name = BODY_FONT
    r.font.size = Pt(24)
    r.font.color.rgb = WHITE
    add_footer(slide, "Spec Coding × 三节点 Agent × Brain Secretary")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide)
    add_title(slide, "这次分享想回答什么问题")
    add_bullets(
        slide,
        [
            "Spec Coding 到底解决了什么问题",
            "为什么单 Agent 在真实项目里往往不够",
            "三节点 Agent 如何形成交付闭环",
            "这套方式在项目里怎么工程化落地",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, TEAL)
    add_title(slide, "什么是 Spec Coding", "先把任务说清楚，再让 AI 去执行")
    add_two_column(
        slide,
        "Spec 的四个核心要素",
        ["目标：这轮到底解决什么问题", "约束：哪些能动，哪些不能动", "验证：改完怎么证明成立", "输出：最后交什么结果"],
        "本质不是长 Prompt",
        ["不是“想到什么问什么”", "不是“让模型多看点上下文”", "而是把任务合同先写清楚", "让后续执行、验收、复盘都可对齐"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, CORAL)
    add_title(slide, "为什么单 Agent 不够")
    add_bullets(
        slide,
        [
            "单 Agent 同时承担理解需求、实现代码、宣布完成三个角色",
            "真实项目里最常见的三个问题：理解偏、越边界、没验证就说完成",
            "所以从“AI 会写”走向“AI 能交付”，关键不是模型更强，而是流程要拆角色",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, NAVY)
    add_title(slide, "三节点 Agent 架构")
    add_flow_box(slide, 0.8, 2.2, 3.3, 1.5, "qq-main", "协调节点\n理解需求、拆任务、控边界、汇总结果", NAVY)
    add_flow_box(slide, 4.95, 2.2, 3.1, 1.5, "brain-secretary-dev", "实施节点\n改代码、改脚本、跑验证", TEAL)
    add_flow_box(slide, 8.9, 2.2, 3.0, 1.5, "brain-secretary-review", "验收节点\n找风险、补问题、提返工", CORAL)
    add_arrow(slide, 4.15, 2.95, 4.85, 2.95)
    add_arrow(slide, 8.15, 2.95, 8.8, 2.95)
    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.9), Inches(4.4), Inches(5.5), Inches(1.2))
    note.fill.solid()
    note.fill.fore_color.rgb = LIGHT
    note.line.color.rgb = RGBColor(220, 226, 232)
    tf = note.text_frame
    tf.text = "扩展节点：auto-evolve-main\n不承接 QQ 主入口，用于自动巡检、自动闭环和无人值守进化。"
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(14)
            run.font.color.rgb = TEXT
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, GOLD)
    add_title(slide, "Spec 在三节点之间如何流转")
    steps = [
        ("需求进入", "用户需求或自动巡检任务进入协调节点"),
        ("生成 Spec", "把目标、约束、验证、输出格式写成任务合同"),
        ("实施开发", "dev 节点按合同执行具体工程改动"),
        ("验收复核", "review 节点独立检查风险、回归和漏测"),
        ("汇总结果", "协调节点收敛结果，保留最终结论"),
        ("只看异常", "默认只把异常和待决策项抛给人"),
    ]
    x = 0.6
    for idx, (title, body) in enumerate(steps):
        add_flow_box(slide, x, 2.2, 1.95, 2.0, f"{idx+1}. {title}", body, [NAVY, TEAL, CORAL, GOLD, SLATE, NAVY][idx])
        if idx < len(steps) - 1:
            add_arrow(slide, x + 1.95, 3.2, x + 2.2, 3.2)
        x += 2.1
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, TEAL)
    add_title(slide, "项目是什么")
    add_two_column(
        slide,
        "项目定位",
        [
            "OpenClaw 多 Agent 工程系统",
            "QQ 入口统一接到 qq-main",
            "dev/review 承担实施与验收",
            "Paperclip 用于协同投影和可视化",
        ],
        "这不是单纯聊天机器人",
        [
            "它关注的是工程交付",
            "不是单轮对话效果",
            "目标是把任务做完、做对、做得可追踪",
            "并让人只在异常处介入",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, SLATE)
    add_title(slide, "项目演进路线")
    phases = [
        ("阶段一", "设计约束与角色拆分"),
        ("阶段二", "项目搭建与基础链路打通"),
        ("阶段三", "功能开发与多 Agent 协同"),
        ("阶段四", "部署、验收与自动闭环"),
    ]
    x = 0.75
    colors = [NAVY, TEAL, CORAL, GOLD]
    for idx, (title, body) in enumerate(phases):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.5), Inches(2.9), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[idx]
        shape.line.color.rgb = colors[idx]
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = TITLE_FONT
        r1.font.bold = True
        r1.font.size = Pt(18)
        r1.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = BODY_FONT
        r2.font.size = Pt(12)
        r2.font.color.rgb = WHITE
        x += 2.9
    hint = slide.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.8), Inches(0.6))
    p = hint.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "如果你后面要讲真实“10 天 / 指令数 / 阶段数据”，这一页可以替换成你自己的项目统计。"
    r.font.name = BODY_FONT
    r.font.size = Pt(14)
    r.font.color.rgb = SLATE
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, CORAL)
    add_title(slide, "典型案例怎么讲")
    add_bullets(
        slide,
        [
            "案例一：AI 驱动需求拆解与设计，重点讲协调节点如何把模糊需求转成结构化任务",
            "案例二：Spec 驱动实施开发，重点讲 dev 节点如何按合同做实现，而不是自由发挥",
            "案例三：Spec 驱动验收与返工，重点讲 review 节点如何发现问题并把任务打回",
            "案例四：复杂问题排障，重点讲协调收敛问题、实施试验修复、验收判断是否成立",
        ],
        top=1.95,
        height=4.8,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, NAVY)
    add_title(slide, "工程化落地：不是聊天，而是系统")
    add_two_column(
        slide,
        "规则与边界",
        [
            "CLAUDE.md / Rules 约束角色和边界",
            "统一运维入口：ops_manager.py",
            "部署真源：deployment_manifest.json",
        ],
        "结果与监督",
        [
            "结构化报告约束最终输出",
            "异常优先汇总：人只看异常",
            "review 证据校验：不靠嘴上说完成",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_top_band(slide, TEAL)
    add_title(slide, "开发者角色的重构")
    add_bullets(
        slide,
        [
            "从“亲手执行每一步”，转向“定义目标、设规则、看异常、做取舍”",
            "人的价值更多体现在目标判断、边界控制、风险接受与否",
            "执行本身越来越可以交给 AI，但验收和责任不能消失",
        ],
        top=2.1,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title(slide, "一句话总结", "Spec 让 AI 有边界，多 Agent 让 AI 有分工，验收机制让 AI 的结果真正可用。", dark=True)
    end_box = slide.shapes.add_textbox(Inches(0.85), Inches(3.1), Inches(11.5), Inches(1.3))
    p = end_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "真正值得期待的，不只是 AI 会不会继续写更多代码，而是我们能不能把 AI 组织进一个更稳定、更可信、更适合交付的工程系统里。"
    r.font.name = BODY_FONT
    r.font.size = Pt(22)
    r.font.color.rgb = WHITE
    add_footer(slide, "Thanks")

    prs.save(str(OUTPUT_PPTX))
    print(OUTPUT_PPTX)


if __name__ == "__main__":
    build_deck()
